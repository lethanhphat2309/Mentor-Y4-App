import streamlit as st
import google.generativeai as genai
import PyPDF2
import json
import re
import random
import gspread 
import requests
import io
from google.oauth2.service_account import Credentials 
from PIL import Image
from pptx import Presentation # MỚI: Thư viện chuyên đọc PowerPoint

# --- 1. CẤU HÌNH GIAO DIỆN & BẢN QUYỀN THÀNH PHÁT ---
st.set_page_config(page_title="Mentor Y4 - Thành Phát", page_icon="👨‍⚕️", layout="centered")

st.markdown("""
    <style>
    .stRadio > label { font-weight: bold; color: #D4AF37; font-size: 16px;}
    .stButton>button { border-radius: 12px; border: 2px solid #D4AF37; color: #D4AF37; background-color: transparent; font-weight: bold; transition: 0.3s;}
    .stButton>button:hover { background-color: #D4AF37; color: white;}
    .footer { position: fixed; left: 0; bottom: 0; width: 100%; background-color: transparent; color: #888; text-align: center; font-size: 13px; padding: 10px; border-top: 1px solid #333;}
    </style>
""", unsafe_allow_html=True)

st.title("👨‍⚕️ Mentor Y Khoa Cá Nhân")
st.caption("☁️ Hệ thống học tập thiết kế độc quyền cho **Thành Phát**")

# --- 2. KẾT NỐI GOOGLE SHEETS BÍ MẬT ---
@st.cache_resource
def init_gsheets():
    try:
        creds_json = json.loads(st.secrets["GOOGLE_CREDENTIALS"])
        scopes = ["https://www.googleapis.com/auth/spreadsheets", "https://www.googleapis.com/auth/drive"]
        creds = Credentials.from_service_account_info(creds_json, scopes=scopes)
        client = gspread.authorize(creds)
        return client.open("Mentor_Y4_Database") 
    except Exception as e:
        return None

sheet_db = init_gsheets()

# --- 3. BỘ NHỚ HỆ THỐNG & TỰ ĐỘNG TẢI DỮ LIỆU ---
def parse_sheet_data(worksheet):
    try:
        records = worksheet.get_all_records()
        data = []
        for row in records:
            options = [opt.strip() for opt in str(row.get("CÁC ĐÁP ÁN", row.get("Các đáp án", ""))).split("\n") if opt.strip()]
            data.append({
                "question": str(row.get("CÂU HỎI", row.get("Câu hỏi", ""))),
                "options": options,
                "answer": str(row.get("ĐÁP ÁN ĐÚNG", row.get("Đáp án đúng", ""))),
                "explanation": str(row.get("GIẢI THÍCH", row.get("Giải thích", "")))
            })
        return data
    except Exception:
        return []

if "data_loaded" not in st.session_state:
    st.session_state.messages = [{"role": "model", "parts": ["Chào Thành Phát! Trợ lý đã trang bị thêm mắt đọc PowerPoint. Gửi slide cho mình nhé!"]}]
    st.session_state.quiz_data = []
    st.session_state.wrong_notebook = []
    st.session_state.current_page = "💬 Chat Mentor"
    st.session_state.last_summary = ""
    st.session_state.data_loaded = True 
    
    if sheet_db:
        try:
            st.session_state.quiz_data = parse_sheet_data(sheet_db.worksheet("QuizBank"))
            st.session_state.wrong_notebook = parse_sheet_data(sheet_db.worksheet("WrongNotebook"))
        except Exception: pass 

def format_for_sheet(data_list):
    rows = [["CÂU HỎI", "CÁC ĐÁP ÁN", "ĐÁP ÁN ĐÚNG", "GIẢI THÍCH"]]
    if not data_list: return rows
    for item in data_list:
        options_str = "\n".join(item['options']) 
        rows.append([item['question'], options_str, item['answer'], item['explanation']])
    return rows

def sync_to_cloud():
    if sheet_db:
        try:
            quiz_ws = sheet_db.worksheet("QuizBank")
            quiz_ws.clear() 
            quiz_ws.update(format_for_sheet(st.session_state.quiz_data)) 
            wrong_ws = sheet_db.worksheet("WrongNotebook")
            wrong_ws.clear()
            wrong_ws.update(format_for_sheet(st.session_state.wrong_notebook))
        except Exception: pass

# --- 4. THANH ĐIỀU HƯỚNG BÊN TRÁI & XỬ LÝ ĐA TỆP ---
with st.sidebar:
    st.markdown("### 👨‍⚕️ Chủ sở hữu: **Thành Phát**")
    st.caption("Phiên bản Y4 - PowerPoint Reader v6.0")
    st.markdown("---")
    
    st.header("⚙️ Hệ Thống")
    api_key = st.text_input("Nhập Google Gemini API Key:", type="password")
    model_choice = st.selectbox("Chọn Bộ Não AI:", ["gemini-3-flash-preview", "gemini-3.0-flash", "gemini-3.0-pro", "gemini-1.5-pro-latest"])
    
    st.markdown("---")
    menu_options = ["💬 Chat Mentor", "📝 Phòng Thi Ảo", "📓 Sổ Tay Câu Sai"]
    selected_page = st.radio("📌 Điều Hướng Ứng Dụng", menu_options, index=menu_options.index(st.session_state.current_page))
    if selected_page != st.session_state.current_page:
        st.session_state.current_page = selected_page
        st.rerun()
    
    pdf_text = ""
    img_data_list = [] 
    
    # KÊNH 1: TẢI FILE TRỰC TIẾP TỪ MÁY (Hỗ trợ thêm .pptx)
    uploaded_files = st.file_uploader("📂 Tải Ảnh/PDF/PPTX từ máy tính", type=["pdf", "txt", "png", "jpg", "jpeg", "pptx"], accept_multiple_files=True)
    if uploaded_files:
        for file in uploaded_files:
            file_ext = file.name.split('.')[-1].lower()
            if file_ext == 'pdf':
                try:
                    reader = PyPDF2.PdfReader(file)
                    extracted_text = ""
                    for page in reader.pages:
                        extracted_text += page.extract_text() or ""
                        if len(extracted_text) > 15000: break
                    pdf_text += f"\n\n[DỮ LIỆU PDF]:\n" + extracted_text[:15000]
                except Exception: pass
            
            # MỚI: Thuật toán đọc PowerPoint
            elif file_ext == 'pptx':
                try:
                    prs = Presentation(file)
                    extracted_text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                extracted_text += shape.text + "\n"
                        if len(extracted_text) > 15000: break # Ngắt an toàn chống treo
                    pdf_text += f"\n\n[DỮ LIỆU SLIDE BÀI GIẢNG]:\n" + extracted_text[:15000]
                except Exception as e: 
                    st.error(f"Lỗi đọc PPTX: {e}")
                    
            elif file_ext in ['png', 'jpg', 'jpeg']:
                try:
                    img = Image.open(file)
                    img_data_list.append(img)
                    st.image(img, caption=f"Đã nạp: {file.name}", use_container_width=True)
                except Exception: pass

    # KÊNH 2: TẢI TỪ GOOGLE DRIVE
    st.markdown("---")
    drive_link = st.text_input("🔗 Hoặc Dán Link Google Drive vào đây:", placeholder="https://drive.google.com/file/d/...")
    
    if drive_link:
        try:
            file_id = None
            match1 = re.search(r"/file/d/([a-zA-Z0-9_-]+)", drive_link)
            match2 = re.search(r"id=([a-zA-Z0-9_-]+)", drive_link)
            if match1: file_id = match1.group(1)
            elif match2: file_id = match2.group(1)
            
            if file_id:
                with st.spinner("☁️ Đang kéo dữ liệu trực tiếp từ Drive..."):
                    url = f"https://drive.google.com/uc?id={file_id}&export=download"
                    response = requests.get(url)
                    
                    if response.status_code == 200:
                        file_bytes = io.BytesIO(response.content)
                        success = False
                        
                        # Thử đọc như PDF
                        try:
                            reader = PyPDF2.PdfReader(file_bytes)
                            extracted_text = ""
                            for page in reader.pages:
                                extracted_text += page.extract_text() or ""
                                if len(extracted_text) > 15000: break
                            pdf_text += f"\n\n[DỮ LIỆU DRIVE - PDF]:\n" + extracted_text[:15000]
                            st.success("✅ Đã kéo xong PDF từ Drive!")
                            success = True
                        except: pass
                        
                        # Thử đọc như PPTX nếu không phải PDF
                        if not success:
                            try:
                                file_bytes.seek(0) # Reset con trỏ file
                                prs = Presentation(file_bytes)
                                extracted_text = ""
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if hasattr(shape, "text"):
                                            extracted_text += shape.text + "\n"
                                    if len(extracted_text) > 15000: break
                                pdf_text += f"\n\n[DỮ LIỆU DRIVE - PPTX]:\n" + extracted_text[:15000]
                                st.success("✅ Đã bóc tách thành công Text từ Slide PPTX trên Drive!")
                                success = True
                            except: pass

                        # Thử đọc như Hình ảnh nếu cả 2 cái trên đều thất bại
                        if not success:
                            try:
                                file_bytes.seek(0)
                                img = Image.open(file_bytes)
                                img_data_list.append(img)
                                st.image(img, caption="Ảnh từ Drive", use_container_width=True)
                                st.success("✅ Đã nạp xong Ảnh từ Drive!")
                                success = True
                            except:
                                st.error("❌ Định dạng không hỗ trợ hoặc file quá nặng. Hãy kiểm tra lại.")
                    else:
                        st.error("❌ Không thể tải! Bạn nhớ bật quyền 'Bất kỳ ai có liên kết' nhé.")
            else:
                st.warning("⚠️ Link chưa đúng định dạng. Hãy copy lại từ Google Drive.")
        except Exception as e:
            st.error(f"Lỗi hệ thống kéo file: {e}")

    st.markdown("---")
    if st.button("🔄 Ép Đồng Bộ Lên Cloud Ngay"):
        sync_to_cloud()
        st.success("Đã lưu an toàn lên Cloud!")

# ==========================================
# CHẾ ĐỘ 1: CHAT MENTOR
# ==========================================
if st.session_state.current_page == "💬 Chat Mentor":
    for msg in st.session_state.messages:
        with st.chat_message("ai" if msg["role"] == "model" else "user"):
            st.markdown(msg["parts"][0], unsafe_allow_html=True)

    user_input = st.chat_input("Nhắn Mentor (VD: Tóm tắt bài này / Trình bày chi tiết bài này)")

    if user_input:
        st.chat_message("user").write(user_input)
        st.session_state.messages.append({"role": "user", "parts": [user_input]})
        
        if api_key:
            genai.configure(api_key=api_key)
            
            # --- LUỒNG A: TẠO TRẮC NGHIỆM ---
            if any(kw in user_input.lower() for kw in ["trắc nghiệm", "câu hỏi", "test", "đề thi"]):
                model = genai.GenerativeModel(model_name=model_choice, generation_config={"response_mime_type": "application/json"})
                schema_instruction = (
                    "Bạn là Mentor Y Khoa. Dựa vào nội dung, hãy tạo MẢNG JSON chứa trắc nghiệm:\n"
                    "[\n"
                    "  {\n"
                    "    \"question\": \"Câu hỏi...\",\n"
                    "    \"options\": [\"A. ...\", \"B. ...\", \"C. ...\", \"D. ...\"],\n"
                    "    \"answer\": \"A. ...\",\n"
                    "    \"explanation\": \"Giải thích chi tiết...\"\n"
                    "  }\n"
                    "]\n"
                    "TUYỆT ĐỐI KHÔNG để dấu phẩy (,) ở cuối phần tử cuối cùng."
                )
                
                prompt_parts = [schema_instruction + pdf_text + "\n\nYêu cầu: " + user_input]
                if img_data_list: prompt_parts.extend(img_data_list)
                    
                with st.spinner("Đang soạn đề thi..."):
                    try:
                        response = model.generate_content(prompt_parts)
                        clean_json = re.sub(r',\s*]', ']', response.text)
                        clean_json = re.sub(r',\s*}', '}', clean_json)
                        new_questions = json.loads(clean_json)
                        
                        st.session_state.quiz_data.extend(new_questions)
                        sync_to_cloud() 
                        
                        st.session_state.messages.append({"role": "model", "parts": [f"Đã nạp thêm {len(new_questions)} câu hỏi vào Ngân Hàng Đề!"]})
                        st.session_state.current_page = "📝 Phòng Thi Ảo"
                        st.rerun() 
                    except Exception as e:
                        st.error(f"Lỗi hệ thống: Quota hoặc AI đang quá tải. Đợi 1 chút nhé!")
            
            # --- LUỒNG B: CHAT & TÓM TẮT ---
            else:
                model = genai.GenerativeModel(model_choice)
                
                format_instruction = (
                    "Bạn là Trưởng khoa Y. Tùy thuộc vào yêu cầu của người dùng, "
                    "hãy đáp ứng chính xác về mặt nội dung. "
                    "BẮT BUỘC phải trình bày chuẩn Markdown: "
                    "1. Dùng Tiêu đề lớn (##) cho các phần chính. "
                    "2. Dùng gạch đầu dòng (-) rõ ràng. "
                    "3. BÔI ĐẬM (**) các từ khóa y khoa, tên thuốc, triệu chứng quan trọng. "
                    "4. Cấu trúc mạch lạc, trực quan."
                )
                
                prompt_parts = [format_instruction + "\n\nNội dung tài liệu:\n" + pdf_text + "\n\nYêu cầu của người dùng: " + user_input]
                if img_data_list: prompt_parts.extend(img_data_list)
                    
                with st.spinner("Mentor đang đọc và phân tích tài liệu theo ý Thành Phát..."):
                    try:
                        response = model.generate_content(prompt_parts)
                        st.chat_message("ai").markdown(response.text)
                        st.session_state.messages.append({"role": "model", "parts": [response.text]})
                        
                        st.session_state.last_summary = response.text
                        st.rerun() 
                    except Exception as e:
                        st.error(f"Lỗi AI: Đợi 1 chút rồi thử lại nhé!")
        else:
            st.warning("Nhớ nhập API Key nhé Thành Phát!")

    if st.session_state.last_summary:
        st.markdown("---")
        st.download_button(
            label="📥 Tải bài Tóm tắt/Giải thích vừa rồi về máy",
            data=st.session_state.last_summary,
            file_name="Giao_An_Y_Khoa_ThanhPhat.md",
            mime="text/markdown",
            help="Bấm để tải file. Mở bằng Google Docs hoặc Notion để giữ nguyên định dạng in đậm và gạch đầu dòng cực đẹp!"
        )

# ==========================================
# CHẾ ĐỘ 2 & 3: GIỮ NGUYÊN
# ==========================================
elif st.session_state.current_page == "📝 Phòng Thi Ảo":
    st.subheader(f"📝 Ngân Hàng Đề Thi Cloud ({len(st.session_state.quiz_data)} câu)")
    if len(st.session_state.quiz_data) == 0:
        st.info("Chưa có câu hỏi. Hãy về tab Chat Mentor yêu cầu tạo trắc nghiệm.")
    else:
        if st.button("🎲 Xáo Trộn Đề (Ôn Tập Ngẫu Nhiên)"):
            random.shuffle(st.session_state.quiz_data)
            st.rerun()
        st.markdown("---")
        for idx, q in enumerate(st.session_state.quiz_data):
            st.markdown(f"**Câu {idx+1}: {q['question']}**")
            choice = st.radio("Chọn đáp án:", q['options'], key=f"radio_{idx}", index=None)
            if st.button(f"Nộp đáp án Câu {idx+1}", key=f"btn_{idx}"):
                if choice is None:
                    st.warning("Chưa chọn đáp án kìa!")
                elif choice == q['answer']:
                    st.success(f"🎉 ĐÚNG RỒI! \n\n**Giải thích sâu:** {q['explanation']}")
                    st.balloons()
                else:
                    st.error(f"❌ SAI RỒI! \n\n**Đáp án đúng:** {q['answer']} \n\n**Giải thích sâu:** {q['explanation']}")
                    if not any(item['question'] == q['question'] for item in st.session_state.wrong_notebook):
                        st.session_state.wrong_notebook.append(q)
                        sync_to_cloud() 
            st.markdown("---")

elif st.session_state.current_page == "📓 Sổ Tay Câu Sai":
    st.subheader("📓 Góc Ôn Tập Của Thành Phát")
    if len(st.session_state.wrong_notebook) == 0:
        st.success("Tuyệt vời! Bạn chưa làm sai câu nào.")
    else:
        st.warning(f"Có {len(st.session_state.wrong_notebook)} câu cần ôn lại:")
        for idx, wq in enumerate(st.session_state.wrong_notebook):
            with st.expander(f"⚠️ {wq['question']}"):
                st.error(f"**Đáp án đúng:** {wq['answer']}")
                st.info(f"**Cơ chế bệnh sinh:** {wq['explanation']}")

# --- 5. CHÂN TRANG BẢN QUYỀN THÀNH PHÁT ---
st.markdown("""
    <div class="footer">
        <p>© 2024 - Bản quyền thuộc về <b>Thành Phát</b> | Phiên bản PowerPoint Reader v6.0</p>
    </div>
""", unsafe_allow_html=True)
